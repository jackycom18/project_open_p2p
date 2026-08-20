# -*- coding: utf-8 -*-
"""生成课题六答辩 PPT（深色科技风，16:9）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- 配色（深色科技风）----
BG = RGBColor(0x0D, 0x11, 0x17)
PANEL = RGBColor(0x16, 0x1B, 0x22)
BORDER = RGBColor(0x30, 0x36, 0x3D)
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
BLUE = RGBColor(0x2F, 0x81, 0xF7)
GOLD = RGBColor(0xD2, 0x99, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0xC9, 0xD1, 0xD9)
MUTED = RGBColor(0x8B, 0x94, 0x9E)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
RED = RGBColor(0xF8, 0x51, 0x49)

FONT = "微软雅黑"
FONT_EN = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, runs, size=16, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         valign=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]
    for i, (t, opt) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opt.get("align", align)
        r = p.add_run()
        r.text = t
        r.font.size = Pt(opt.get("size", size))
        r.font.color.rgb = opt.get("color", color)
        r.font.bold = opt.get("bold", bold)
        r.font.name = opt.get("font", font)
    tf.vertical_anchor = valign
    return tb


def header(slide, title, idx=None):
    """内容页统一页头：左侧竖色条 + 标题 + 右上角页码小标签。"""
    rect(slide, 0.6, 0.55, 0.12, 0.6, BLUE)
    text(slide, 0.9, 0.5, 11.0, 0.7, title, size=30, color=WHITE, bold=True)
    if idx:
        text(slide, 12.3, 0.55, 0.6, 0.4, idx, size=12, color=MUTED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body_runs, title_color=BLUE):
    rect(slide, x, y, w, h, PANEL, line=BORDER)
    rect(slide, x, y, 0.08, h, title_color)
    text(slide, x + 0.25, y + 0.18, w - 0.4, 0.4, title, size=15, color=title_color, bold=True)
    text(slide, x + 0.25, y + 0.65, w - 0.5, h - 0.8, body_runs, size=12.5, color=TEXT)


def big_stat(slide, x, y, num, label, num_color=BLUE, w=3.0):
    text(slide, x, y, w, 1.0, num, size=48, color=num_color, bold=True, align=PP_ALIGN.CENTER)
    text(slide, x, y + 1.0, w, 0.7, label, size=13, color=MUTED, align=PP_ALIGN.CENTER)


# ================= 第 1 页：封面 =================
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 13.333, 7.5, BG)
# 顶部与底部装饰色条
rect(s, 0, 0, 13.333, 0.18, BLUE)
rect(s, 0, 7.32, 13.333, 0.18, BLUE)
# 左侧大色块（装饰）
rect(s, 0.8, 1.7, 0.16, 3.4, GOLD)

text(s, 1.3, 1.7, 10.5, 0.5, "课题六 · 行为克隆游戏智能体", size=18, color=GOLD, bold=True)
text(s, 1.3, 2.35, 11.0, 1.4, "让 AI 像人类一样\n看着画面打游戏", size=44, color=WHITE, bold=True)
text(s, 1.3, 3.85, 11.0, 0.5, "Open P2P（Pixels2Play）· 文本指令 + 推理性能", size=20, color=ICE)

text(s, 1.3, 5.3, 6.0, 0.5, [
    ("答辩人：向子坚", {"size": 16, "color": TEXT, "bold": True}),
], size=16)
text(s, 1.3, 5.85, 8.0, 0.4, "腾讯 IEG 课题 · AI 开发实践", size=13, color=MUTED)

# ================= 第 2 页：课题与使用场景 =================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "课题与使用场景", "02")

card(s, 0.6, 1.5, 6.0, 2.2, "选题理由", [
    ("行为克隆方法成熟、可复现；150M 档模型单卡 8–16GB 显存即可推理，契合课程算力。", {}),
    ("官方架构原生支持文本指令（--input_text），与扩展方向契合，研究空间大。", {}),
])
card(s, 6.8, 1.5, 6.0, 2.2, "使用场景", [
    ("为游戏测试 / 内容验证人员提供可被文本指令指挥的 AI 游戏操作智能体。", {}),
    ("输入画面 + 一句指令 → 输出符合意图的键鼠动作。", {}),
])
card(s, 0.6, 3.95, 12.2, 2.0, "课内交付边界", [
    ("无真机外设：不依赖真实游戏环境与 Windows Recap 双卡方案；", {"bold": True}),
    ("算力受限：仅推理与评测，不做微调；", {"bold": True}),
    ("交付形态：以离线推理 + 评测 + 对比录屏为主，不承诺实时端到端 <50ms 链路。", {"bold": True}),
])

# ================= 第 3 页：必做与验收（5 项 MVP）=================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "必做与验收 · 5 项 MVP", "03")

rows = [
    ("M1", "跑通官方 150M 推理", "README 一键复现推理输出"),
    ("M2", "评测脚本 + 测试集", "测试集 ≤200 帧，输出按键准确率 + 鼠标相关"),
    ("M3", "未微调基线", "按键准确率 ≥55%，鼠标 r ≥0.5"),
    ("M4", "对比录屏", "同一片段模型输出 vs 人类标注"),
    ("M5", "归档 + 报告", "代码归档 + 指标表 + ≥3000 字报告"),
]
y = 1.6
for i, (m, t, d) in enumerate(rows):
    rect(s, 0.6, y, 12.2, 0.92, PANEL, line=BORDER)
    rect(s, 0.6, y, 0.9, 0.92, NAVY)
    text(s, 0.6, y + 0.22, 0.9, 0.5, m, size=18, color=ICE, bold=True, align=PP_ALIGN.CENTER)
    text(s, 1.75, y + 0.13, 4.3, 0.65, t, size=16, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    text(s, 6.0, y + 0.13, 6.6, 0.65, d, size=13, color=MUTED, valign=MSO_ANCHOR.MIDDLE)
    y += 1.05

# ================= 第 4 页：方案组成与技术选型 =================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "方案组成与技术选型", "04")

card(s, 0.6, 1.5, 6.0, 2.5, "模型（官方复用）", [
    ("Open P2P 150M 预训练权重 + Gemma 文本分词器", {"bold": True}),
    ("选择理由：架构原生支持文本指令，官方文档与社区完整", {}),
])
card(s, 6.8, 1.5, 6.0, 2.5, "评测链路（自研）", [
    ("eval/ 八件套：客户端 / 选测试集 / 评测 / 指令对照 / 性能评测 / 汇总 / 录屏 / Web", {"bold": True}),
    ("选择理由：复用官方 UDS 通信协议，指标可量化、可复现", {}),
])
card(s, 0.6, 4.2, 12.2, 1.6, "部署环境", [
    ("Linux 服务器（NVIDIA RTX 5880 Ada 48GB）· uv 管理依赖 · HF 离线模式加载 Gemma", {"bold": True}),
    ("选择理由：推理需 NVIDIA GPU，Linux 驱动环境比 Windows 原生更稳", {}),
])

# ================= 第 5 页：主路径架构 =================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "主路径架构", "05")

steps = [
    ("select_testset.py", "冻结测试集\n≤200 帧"),
    ("inference.py", "推理服务器\nUDS 通信"),
    ("p2p_client", "发帧\n收动作"),
    ("evaluate.py", "读标注+输出\n算指标"),
    ("summarize/record", "汇总表\n对比录屏"),
]
x = 0.6
for i, (t, d) in enumerate(steps):
    rect(s, x, 2.6, 2.1, 1.6, PANEL, line=BORDER)
    text(s, x, 2.8, 2.1, 0.5, t, size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, 3.35, 2.1, 0.7, d, size=12, color=TEXT, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        text(s, x + 2.05, 3.15, 0.45, 0.5, "→", size=20, color=MUTED, align=PP_ALIGN.CENTER)
    x += 2.45

text(s, 0.6, 4.7, 12.2, 1.4, [
    ("核心机制：画面帧（192×192）→ 模型 → Action（按键 + 鼠标位移）", {"bold": True, "color": ICE}),
    ("文本指令通过服务器 stdin 注入，运行中逐条生效；推理性能通过 KV Cache / 全量重算切换对比。", {"color": TEXT}),
], size=14)

# ================= 第 6 页：实验环境 =================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "实验环境与部署", "06")

card(s, 0.6, 1.6, 5.8, 2.4, "硬件", [
    ("GPU：NVIDIA RTX 5880 Ada（48GB）", {"bold": True}),
    ("CPU / 内存：服务器规格", {}),
    ("网络：Tailscale 组网远程访问", {}),
])
card(s, 6.6, 1.6, 6.2, 2.4, "软件", [
    ("OS：Linux（Ubuntu）+ uv 依赖管理", {"bold": True}),
    ("模型：150M checkpoint + Gemma tokenizer", {}),
    ("数据：p2p-toy-examples（3 片段）", {}),
])
card(s, 0.6, 4.25, 12.2, 1.5, "复现要点", [
    ("① setup_wsl.sh 一键搭环境 → ② download 权重/数据 → ③ inference.py 启动 → ④ eval/ 评测", {"bold": True}),
    ("HF 离线模式（HF_HUB_OFFLINE=1）解决 Tailscale 下 SentenceTransformer 联网卡顿。", {}),
])

# ================= 第 7 页：基线评测结果（M3）=================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "基线评测结果（M3）", "07")

big_stat(s, 1.0, 2.2, "48.50%", "按键准确率（参考 ≥55%）", GOLD)
big_stat(s, 5.0, 2.2, "0.057", "鼠标相关系数（参考 ≥0.5）", GOLD)
big_stat(s, 9.0, 2.2, "14.42", "位移 MAE（报告项）", BLUE)

card(s, 0.6, 4.2, 12.2, 2.0, "结果解读", [
    ("按键 48.50%、鼠标相关 0.057，未达课内参考线，属预期范围内。", {"bold": True}),
    ("官方论文未采用鼠标相关系数（因鼠标困惑度 noisy）；本课指标为课内自定义。", {}),
    ("150M 未微调 + toy 跨域数据 + 静止开局帧，是未达标的主因。", {}),
])

# ================= 第 8 页：文本指令扩展（B）=================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "扩展 B：文本指令", "08")

big_stat(s, 0.6, 1.9, "100%", "带指令完成率\n（片段A）", GREEN, w=2.95)
big_stat(s, 3.6, 1.9, "0%", "不带指令完成率\n（片段A）", RED, w=2.95)
big_stat(s, 6.6, 1.9, "+31pp", "按键准确率增益\n（片段A）", GREEN, w=2.95)
big_stat(s, 9.6, 1.9, "+10pp", "按键准确率增益\n（片段B）", BLUE, w=2.95)

card(s, 0.6, 4.0, 12.2, 2.0, "结论", [
    ("指令\"持续向左移动\"使任务完成率从 0% 升至 100%，按键准确率 24.5%→55.5%。", {"bold": True}),
    ("两个片段按键准确率均正向增益（+31pp / +10pp），验证\"文本指令能辅助模型执行意图明确的操作\"。", {}),
])

# ================= 第 9 页：推理性能扩展（C）=================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "扩展 C：推理性能", "09")

big_stat(s, 0.9, 1.9, "15.26x", "KV Cache\n加速比", GREEN)
big_stat(s, 4.3, 1.9, "75.39", "KV Cache\n帧率 fps", BLUE)
big_stat(s, 7.7, 1.9, "3.94x", "torch.compile\n加速比", GOLD)

rows = [
    ("配置", "耗时(s)", "帧率(fps)", "按键一致率"),
    ("KV Cache（优化后）", "2.653", "75.39", "44.5%"),
    ("全量重算（优化前）", "40.470", "4.94", "23.0%"),
    ("关闭编译（对照组）", "10.465", "19.11", "35.0%"),
]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.6), Inches(4.1), Inches(12.2), Inches(1.9)).table
tbl.columns[0].width = Inches(3.6)
tbl.columns[1].width = Inches(2.8)
tbl.columns[2].width = Inches(2.8)
tbl.columns[3].width = Inches(3.0)
for r in range(len(rows)):
    for c in range(4):
        cell = tbl.cell(r, c)
        cell.text = rows[r][c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else PANEL
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.name = FONT
        p.runs[0].font.color.rgb = WHITE if r == 0 else TEXT
        p.runs[0].font.bold = (r == 0)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# ================= 第 10 页：软件系统展示 =================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "软件系统展示（Web 演示）", "10")

card(s, 0.6, 1.6, 6.0, 4.4, "系统功能", [
    ("① 启动 / 停止推理服务器（支持指令模式）", {"bold": True}),
    ("② 选择测试片段，逐帧推理并播放", {"bold": True}),
    ("③ 输入文本指令，观察动作变化", {"bold": True}),
    ("④ 实时对比：模型预测 vs 人类标注", {"bold": True}),
    ("⑤ 服务器日志实时查看", {"bold": True}),
])
card(s, 6.8, 1.6, 6.0, 4.4, "展示流程", [
    ("Step 1：浏览器打开 Web 界面", {"bold": True}),
    ("Step 2：启动推理服务器（1-2 分钟加载）", {"bold": True}),
    ("Step 3：选片段 → 推理 → 逐帧播放", {"bold": True}),
    ("Step 4：发送指令 → 观察动作变化", {"bold": True}),
    ("（下方留白处可粘贴实际运行截图）", {"color": MUTED}),
])

# ================= 第 11 页：不足与展望 =================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "不足与展望", "11")

card(s, 0.6, 1.6, 6.0, 4.4, "不足", [
    ("未做微调（算力受限），基线未达参考线", {}),
    ("未做实时游戏推理（Windows Recap + 双卡）", {}),
    ("指令片段数量较少（仅 2 个）", {}),
    ("未纳入官方口径 RCE / 程序化环境评测", {}),
])
card(s, 6.8, 1.6, 6.0, 4.4, "展望", [
    ("微调（扩展 A）：按键 +8pp 或鼠标 +0.08", {}),
    ("换 600M / 1200M 更大模型（scaling law）", {}),
    ("采样策略 + 温度网格搜索（零成本）", {}),
    ("补齐官方键盘困惑度 RCE 指标", {}),
])

# ================= 第 12 页：总结 =================
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 0, 13.333, 0.18, GOLD)
rect(s, 0, 7.32, 13.333, 0.18, GOLD)

text(s, 0.8, 1.7, 11.7, 0.5, "总结", size=20, color=GOLD, bold=True)
text(s, 0.8, 2.5, 11.7, 2.4, [
    ("完整跑通 Open P2P 150M 行为克隆模型：环境 → 推理 → 评测 → 演示全链路可复现；", {"size": 20, "color": WHITE, "bold": True}),
    ("自研评测工具链产出按键准确率、鼠标相关系数、任务完成率、推理性能四类可量化指标；", {"size": 20, "color": WHITE, "bold": True}),
    ("完成文本指令（B）+ 推理性能（C）两个扩展，均取得显著正向结论。", {"size": 20, "color": WHITE, "bold": True}),
])
text(s, 0.8, 5.6, 11.7, 0.5, "谢谢观看 · 欢迎提问", size=22, color=ICE, bold=True, align=PP_ALIGN.CENTER)

prs.save(r"C:\Users\Xzj13\CodeBuddy\project\课题六_行为克隆游戏智能体_答辩.pptx")
print("PPT 已生成")
